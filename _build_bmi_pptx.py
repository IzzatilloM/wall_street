# -*- coding: utf-8 -*-
"""Wall Street CRM — BMI himoyasi uchun 17 slaydli prezentatsiya."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

SHOTS = os.path.join(os.path.dirname(__file__), "_pptx_shots")
OUT = os.path.join(os.path.dirname(__file__), "Wall_Street_BMI_Prezentatsiya.pptx")

# ---------- palette ----------
BG    = RGBColor(0x0B, 0x2A, 0x2C)
BG2   = RGBColor(0x0E, 0x35, 0x38)
CARD  = RGBColor(0x10, 0x3E, 0x42)
CARD2 = RGBColor(0x14, 0x4A, 0x4E)
TEAL  = RGBColor(0x0E, 0x5C, 0x60)
MINT  = RGBColor(0x2E, 0xC4, 0xB6)
GOLD  = RGBColor(0xE0, 0xB8, 0x4C)
RED   = RGBColor(0xE0, 0x6A, 0x5C)
TXT   = RGBColor(0xEC, 0xF4, 0xF3)
MUT   = RGBColor(0x9C, 0xB7, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HFONT = "Trebuchet MS"
BFONT = "Calibri"

SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=BG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def _set_font(run, size, color, bold=False, italic=False, font=BFONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line=None, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line:
            p.line_spacing = line
        for seg in para:
            t, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            ital = seg[4] if len(seg) > 4 else False
            fnt  = seg[5] if len(seg) > 5 else BFONT
            r = p.add_run(); r.text = t
            _set_font(r, size, color, bold, ital, fnt)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'),
              {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val':'000000'})
        alpha = clr.makeelement(qn('a:alpha'), {'val':'45000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def img_fit(path, box_w, box_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    bar = box_w / box_h
    if ar > bar:
        w = box_w; h = box_w / ar
    else:
        h = box_h; w = box_h * ar
    return w, h

def framed_image(s, path, cx, cy, box_w, box_h, border=GOLD, bw=1.5):
    w, h = img_fit(path, box_w, box_h)
    x = cx + (box_w - w) / 2
    y = cy + (box_h - h) / 2
    rect(s, x, y, w, h, fill=RGBColor(0,0,0), shadow=True)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    rect(s, x, y, w, h, fill=None, line=border, line_w=bw)
    return x, y, w, h

def num_badge(s, x, y, n, d=0.62):
    rect(s, x, y, d, d, fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, x, y-0.02, d, d, [[(str(n), 24, BG, True, False, HFONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def corner_dots(s):
    rect(s, SW-0.55, 0.32, 0.16, 0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    rect(s, SW-0.32, 0.32, 0.16, 0.16, fill=GOLD, shape=MSO_SHAPE.OVAL)

def bullets(s, x, y, w, items, gap=0.62, dot=GOLD, head_size=13.5, body_size=11):
    cy = y
    for head, body in items:
        rect(s, x, cy+0.07, 0.14, 0.14, fill=dot, shape=MSO_SHAPE.OVAL)
        text(s, x+0.30, cy, w-0.30,
             0.9, [[(head+"  ", head_size, TXT, True, False, HFONT),
                    (body, body_size, MUT, False, False, BFONT)]],
             line=1.04, space_after=0)
        cy += gap
    return cy

def content_slide(n, kicker, title, intro, items, img, img_side="right", accent=GOLD):
    s = slide(); bg(s, BG if n % 2 else BG2)
    corner_dots(s)
    text_w = 4.7
    if img_side == "right":
        tx = 0.7; img_x = 5.75
    else:
        tx = SW - text_w - 0.7; img_x = 0.7
    img_w = 6.85
    num_badge(s, tx, 0.5, n)
    text(s, tx+0.85, 0.52, text_w-0.5, 0.4,
         [[(kicker.upper(), 12, MINT, True, False, HFONT)]], space_after=0)
    text(s, tx, 1.32, text_w, 1.2,
         [[(title, 26, TXT, True, False, HFONT)]], line=0.98, space_after=0)
    text(s, tx, 2.78, text_w, 1.2,
         [[(intro, 12.5, MUT, False, False, BFONT)]], line=1.12, space_after=0)
    bullets(s, tx, 4.05, text_w, items)
    framed_image(s, img, img_x, 1.15, img_w, 5.2, border=accent, bw=1.5)
    return s

P = lambda name: os.path.join(SHOTS, name)

# ============================================================
# SLIDE 1 — TITUL
# ============================================================
s = slide(); bg(s, BG)
rect(s, 9.4, -1.8, 6.0, 6.0, fill=BG2, shape=MSO_SHAPE.OVAL)
rect(s, 10.8, 3.6, 4.4, 4.4, fill=CARD, shape=MSO_SHAPE.OVAL)
# logo badge
rect(s, 0.9, 0.8, 1.25, 1.25, fill=TEAL, shape=MSO_SHAPE.OVAL, line=GOLD, line_w=2.0)
text(s, 0.9, 0.85, 1.25, 1.15, [[("WS", 30, GOLD, True, False, HFONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 2.35, 0.95, 8, 0.5, [[("WALL STREET", 22, TXT, True, False, HFONT)]], space_after=0)
text(s, 2.36, 1.48, 8, 0.4, [[("EDUCATION CENTRE", 13, MINT, True, False, HFONT)]], space_after=0)

text(s, 0.9, 2.45, 7.6, 0.5,
     [[("BITIRUV MALAKAVIY ISHI", 14, GOLD, True, False, HFONT)]], space_after=0)
text(s, 0.9, 3.05, 9.7, 2.0,
     [[("O'quv markazi faoliyatini boshqarishning", 30, TXT, True, False, HFONT)],
      [("raqamli axborot tizimini ishlab chiqish", 30, TXT, True, False, HFONT)]],
     line=1.04, space_after=2)
text(s, 0.93, 4.62, 9.4, 0.6,
     [[("Wall Street o'quv markazi misolida — Web CRM, mobil ilova, "
        "Telegram bot va REST API", 14, MUT, False, True, BFONT)]], line=1.18)

# author / supervisor block
rect(s, 0.9, 5.55, 5.7, 1.35, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, line=TEAL, line_w=1.0)
text(s, 1.2, 5.75, 5.2, 0.4, [[("Bajardi:  ", 12, MINT, True, False, HFONT),
                               ("Izzatillo Musayev", 13, TXT, True, False, HFONT)]], space_after=4)
text(s, 1.2, 6.18, 5.2, 0.4, [[("Ilmiy rahbar:  ", 12, MINT, True, False, HFONT),
                               ("____________________", 12, MUT, False, False, BFONT)]], space_after=0)
text(s, 1.2, 6.55, 5.2, 0.4, [[("Toshkent — 2026", 11, MUT, False, True, BFONT)]], space_after=0)

# ============================================================
# SLIDE 2 — DOLZARBLIK (Muammo vs Yechim)
# ============================================================
s = slide(); bg(s, BG2); corner_dots(s)
num_badge(s, 0.7, 0.55, 2)
text(s, 1.55, 0.56, 9, 0.4, [[("MAVZUNING DOLZARBLIGI", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 12, 0.7, [[("Nega bu tizim kerak?", 30, TXT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 2.0, 12, 0.6,
     [[("Bugungi kunda ko'plab o'quv markazlari talaba, to'lov va davomatni qog'oz "
        "daftar yoki Excelda yuritadi. Bu xatolar, ma'lumot yo'qolishi va vaqt "
        "isrofiga olib keladi.", 13, MUT, False, False, BFONT)]], line=1.16)

colw = 5.85
# Muammo card
rect(s, 0.7, 3.05, colw, 3.65, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, line=RED, line_w=1.2)
rect(s, 0.7, 3.05, 0.12, 3.65, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
text(s, 1.1, 3.3, colw-0.6, 0.5, [[("An'anaviy usul — muammolar", 16, RED, True, False, HFONT)]], space_after=0)
probs = ["Ma'lumotlar daftar va Excelda tarqoq",
         "To'lov va qarzlarni qo'lda hisoblash, xatolar",
         "Davomat va statistikani kuzatib bo'lmaydi",
         "Chek va hisobotlar qo'lda tayyorlanadi",
         "Talaba o'z holatini ko'ra olmaydi"]
cy = 3.95
for p in probs:
    text(s, 1.1, cy, colw-0.7, 0.5, [[("✕  ", 12, RED, True, False, BFONT),
                                      (p, 12, TXT, False, False, BFONT)]], line=1.05, space_after=0)
    cy += 0.52
# Yechim card
x2 = 0.7 + colw + 0.33
rect(s, x2, 3.05, colw, 3.65, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, line=MINT, line_w=1.2)
rect(s, x2, 3.05, 0.12, 3.65, fill=MINT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
text(s, x2+0.4, 3.3, colw-0.6, 0.5, [[("Raqamli tizim — yechim", 16, MINT, True, False, HFONT)]], space_after=0)
sols = ["Barcha ma'lumot yagona bazada",
        "To'lov, qarz va foiz avtomatik hisoblanadi",
        "Davomat va daromad jonli statistikada",
        "PDF chek — QR va pechat bilan bir tugmada",
        "Talaba mobil ilovada hammasini ko'radi"]
cy = 3.95
for p in sols:
    text(s, x2+0.4, cy, colw-0.7, 0.5, [[("✓  ", 12, MINT, True, False, BFONT),
                                        (p, 12, TXT, False, False, BFONT)]], line=1.05, space_after=0)
    cy += 0.52

# ============================================================
# SLIDE 3 — MAQSAD VA VAZIFALAR
# ============================================================
s = slide(); bg(s, BG); corner_dots(s)
num_badge(s, 0.7, 0.55, 3)
text(s, 1.55, 0.56, 9, 0.4, [[("MAQSAD VA VAZIFALAR", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 12, 0.7, [[("Ishning maqsadi va vazifalari", 30, TXT, True, False, HFONT)]], space_after=0)
# maqsad banner
rect(s, 0.7, 2.05, 11.93, 1.15, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, line=GOLD, line_w=1.2)
text(s, 1.05, 2.24, 1.6, 0.7, [[("MAQSAD", 12, GOLD, True, False, HFONT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, 2.5, 2.24, 9.9, 0.75,
     [[("O'quv markazining talaba, kurs, davomat, to'lov va hisobot jarayonlarini "
        "yagona avtomatlashtirilgan tizimda birlashtirish.", 14, TXT, False, False, BFONT)]],
     anchor=MSO_ANCHOR.MIDDLE, line=1.12, space_after=0)
# vazifalar 2x3 grid
tasks = [
    ("Talaba va kurs boshqaruvi", "ro'yxatga olish, profil, qidiruv va filtr."),
    ("To'lov tizimi", "to'lov qabul qilish, qarz nazorati, PDF chek."),
    ("Davomat va coin", "darslarni belgilash, rag'batlantirish balli."),
    ("Hisobot va tahlil", "daromad, statistika va eksport."),
    ("Mobil ilova + API", "talaba uchun Flutter ilova, JWT xavfsizlik."),
    ("Telegram bot", "avtomatik ro'yxatga olish va lead yig'ish."),
]
gw, gh = 3.85, 1.55
gx0, gy0 = 0.7, 3.55
gapx, gapy = 0.19, 0.2
for i, (h, b) in enumerate(tasks):
    r, c = divmod(i, 3)
    px = gx0 + c*(gw+gapx)
    py = gy0 + r*(gh+gapy)
    rect(s, px, py, gw, gh, fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, line=TEAL, line_w=1.0)
    rect(s, px+0.28, py+0.28, 0.5, 0.5, fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, px+0.28, py+0.27, 0.5, 0.5, [[(str(i+1), 18, BG, True, False, HFONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, px+0.95, py+0.27, gw-1.15, 0.5, [[(h, 14.5, TXT, True, False, HFONT)]], line=0.98, space_after=0)
    text(s, px+0.28, py+0.92, gw-0.55, 0.55, [[(b, 11.5, MUT, False, False, BFONT)]], line=1.08, space_after=0)

# ============================================================
# SLIDE 4 — TIZIM ARXITEKTURASI
# ============================================================
s = slide(); bg(s, BG2); corner_dots(s)
num_badge(s, 0.7, 0.55, 4)
text(s, 1.55, 0.56, 11, 0.5, [[("TIZIM ARXITEKTURASI", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 12, 0.8, [[("To'rtta bog'langan komponent, bitta baza", 28, TXT, True, False, HFONT)]], space_after=0)
comps = [
    ("Web CRM — Django", "Administrator va o'qituvchilar uchun to'liq boshqaruv paneli: "
     "talabalar, kurslar, davomat, to'lovlar va hisobotlar.", GOLD),
    ("Mobil ilova — Flutter", "Talaba o'z profili, kurslari, to'lovlari va coin balansini "
     "telefonida ko'radi. JWT token orqali xavfsiz kirish.", MINT),
    ("Telegram bot — aiogram", "Yangi talabalarni ro'yxatga oladi va lead yig'adi — "
     "webhook orqali sayt bilan bitta jarayonda ishlaydi.", GOLD),
    ("REST API — DRF + JWT", "Mobil ilova va saytni bog'lovchi xavfsiz interfeys: "
     "auth, profil, kurslar, to'lovlar va yangiliklar.", MINT),
]
gw, gh = 5.85, 1.95
positions = [(0.7,2.15),(6.78,2.15),(0.7,4.3),(6.78,4.3)]
for (px,py),(h,b,acc) in zip(positions, comps):
    rect(s, px, py, gw, gh, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, line=TEAL, line_w=1.0)
    rect(s, px, py, 0.12, gh, fill=acc, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, px+0.4, py+0.3, gw-0.7, 0.5, [[(h, 17, TXT, True, False, HFONT)]], space_after=0)
    text(s, px+0.4, py+0.85, gw-0.7, 1.0, [[(b, 12, MUT, False, False, BFONT)]], line=1.12)

# ============================================================
# SLIDES 5..13 — SCREENSHOTS
# ============================================================
content_slide(5, "Xavfsiz kirish", "Kirish va rolga asoslangan huquq",
    "Tizimga username yoki Gmail orqali kiriladi. Har bir foydalanuvchi roli "
    "(admin / o'qituvchi / talaba) bo'yicha o'z paneliga yo'naltiriladi.",
    [("Rolga asoslangan kirish", "— admin to'liq, o'qituvchi cheklangan huquq."),
     ("Parolni ko'rsatish", "va “meni eslab qol” imkoniyati."),
     ("Parolni tiklash", "— SMS / Telegram kod orqali."),
     ("Telegram bot", "avtomatik login/parol beradi.")],
    P("01_login.png"), "right")

content_slide(6, "Boshqaruv paneli", "Dashboard — bir qarashda butun markaz",
    "Bosh sahifada markazning asosiy ko'rsatkichlari jonli chiqadi: oylik daromad, "
    "kutilayotgan to'lovlar, faol talabalar va kurslar soni.",
    [("Moliyaviy kartalar", "— oylik daromad va qarzlar."),
     ("Statistika", "— talaba, kurs, o'qituvchi soni."),
     ("Davomat xulosasi", "— bugungi darslar holati."),
     ("Tez navigatsiya", "— barcha bo'limlarga menyudan.")],
    P("03_dashboard.png"), "left")

content_slide(7, "Talabalar", "Talabalar ro'yxati va boshqaruvi",
    "Barcha talabalar yagona ro'yxatda. Qidiruv, filtrlash, yangi talaba qo'shish va "
    "har bir talabaning holatini boshqarish.",
    [("Qidiruv va filtr", "— ism, ID yoki holat bo'yicha."),
     ("Yangi talaba", "qo'shish va tahrirlash."),
     ("Holat almashtirish", "— faol / nofaol qilish."),
     ("Profilga o'tish", "— bitta bosishda to'liq ma'lumot.")],
    P("04_students.png"), "right")

content_slide(8, "Talaba profili", "Har bir talaba — to'liq kartochka",
    "Profilda shaxsiy ma'lumotlar, yozilgan kurslar, to'lov tarixi, davomat foizi va "
    "coin balansi birgalikda ko'rinadi.",
    [("Shaxsiy ma'lumot", "— telefon, ota-ona, manzil, rasm."),
     ("Kurslar va to'lovlar", "bitta joyda."),
     ("Davomat foizi", "har bir kurs bo'yicha."),
     ("Coin balansi", "— rag'batlantirish ko'rsatkichi.")],
    P("05_student_detail.png"), "left")

content_slide(9, "Kurslar", "Kurslar va talabani ko'chirish",
    "Markazdagi barcha kurslar daraja va narx bilan. Kurs yaratish, tahrirlash hamda "
    "talabani bir kursdan boshqasiga bir bosishda ko'chirish.",
    [("Kurslar katalogi", "— daraja, narx, davomiylik."),
     ("Modal oynada yaratish", "— sahifani tark etmasdan."),
     ("Transfer imkoniyati", "— talabani boshqa kursga o'tkazish."),
     ("Kurs ichidagilar", "— talabalar ro'yxati.")],
    P("06b_transfer_modal.png"), "right")

content_slide(10, "Davomat va coin", "Davomat va coin rag'batlantirish tizimi",
    "O'qituvchi har bir dars uchun davomat belgilaydi (keldi/kelmadi/kechikdi), tizim "
    "foizni avtomatik hisoblaydi. Shu oynadan talabaga coin beriladi.",
    [("Davomat belgilash", "— present / absent / late."),
     ("Avtomatik foiz", "har bir talaba bo'yicha."),
     ("Coin berish / ayirish", "— sababi bilan, tarix saqlanadi."),
     ("CSV eksport", "— hisobot uchun.")],
    P("09_coin_award.png"), "left")

content_slide(11, "To'lovlar", "To'lovlar ro'yxati va nazorati",
    "Barcha to'lovlar yagona ro'yxatda: summa, holat, oy va talaba. Yangi to'lovda "
    "narx avtomatik chiqadi, usul (naqd / karta / o'tkazma) tanlanadi.",
    [("To'lovlar ro'yxati", "— summa va holat."),
     ("Filtr va qidiruv", "— oy, talaba bo'yicha."),
     ("Holatlar", "— to'langan / qisman / qarz."),
     ("Avtomatik narx", "— kurs bo'yicha hisoblanadi.")],
    P("10_payments.png"), "right")

content_slide(12, "PDF chek", "To'lov cheki — PDF, QR va pechat",
    "Har bir to'lov uchun rasmiy raqamli kvitansiya: talaba, kurs, o'qituvchi va summa. "
    "Bir tugma bilan QR-kod va markaz pechati bosilgan PDF yuklab olinadi.",
    [("To'liq tafsilot", "— talaba, kurs, oy, summa."),
     ("PDF chek", "— yuklab olish / chop etish."),
     ("QR-kod", "— chekni tekshirish uchun."),
     ("Pechat va raqam", "— rasmiy kvitansiya.")],
    P("12_payment_detail.png"), "left")

content_slide(13, "Hisobotlar", "Reports — moliyaviy tahlil va eksport",
    "Markaz faoliyatining umumiy tahlili grafiklarda: daromad dinamikasi, to'lovlar, "
    "talabalar o'sishi. Maosh va hisobotlar Excel/PDF ga eksport qilinadi.",
    [("Daromad tahlili", "— davr bo'yicha grafiklar."),
     ("To'lov statistikasi", "— to'langan / qarz."),
     ("Xodimlar maoshi", "— hisob va eksport."),
     ("Excel / PDF eksport", "— rasmiy hujjat uchun.")],
    P("13_reports.png"), "right")

# ============================================================
# SLIDE 14 — TELEGRAM BOT
# ============================================================
s = slide(); bg(s, BG); corner_dots(s)
num_badge(s, 0.7, 0.55, 14)
text(s, 1.55, 0.56, 9, 0.4, [[("TELEGRAM BOT", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 7, 1.0, [[("Bot orqali ro'yxatga olish", 30, TXT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 2.15, 6.6, 1.2,
     [[("Telegram bot yangi talabalarni avtomatik ro'yxatga oladi va lead yig'adi. "
        "Webhook orqali ishlaydi — sayt bilan bitta jarayon, bitta baza.", 13.5, MUT, False, False, BFONT)]],
     line=1.18)
bot_items = [
    ("/start menyusi", "— kurslar, narxlar, manzil."),
    ("Telefon va ism", "so'raydi, holatni bazada saqlaydi."),
    ("Avtomatik login", "— kirish ma'lumotlarini beradi."),
    ("Lead yig'ish", "— qiziqqanlar CRM ga tushadi."),
    ("Inline tugmalar", "— kurslarni tanlash uchun."),
]
bullets(s, 0.7, 3.45, 6.6, bot_items, gap=0.6, head_size=13.5, body_size=11.5)
# phone mock
pmx, pmy, pmw, pmh = 8.0, 1.15, 4.4, 5.9
rect(s, pmx, pmy, pmw, pmh, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, line=GOLD, line_w=1.5, shadow=True)
rect(s, pmx, pmy, pmw, 0.95, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
rect(s, pmx, pmy+0.45, pmw, 0.5, fill=TEAL)
rect(s, pmx+0.28, pmy+0.22, 0.5, 0.5, fill=GOLD, shape=MSO_SHAPE.OVAL)
text(s, pmx+0.30, pmy+0.27, 0.5, 0.4, [[("WS", 12, BG, True, False, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, pmx+0.95, pmy+0.20, 3, 0.3, [[("Wall Street bot", 13, WHITE, True, False, HFONT)]], space_after=0)
text(s, pmx+0.95, pmy+0.52, 3, 0.3, [[("online", 10, RGBColor(0xCD,0xEA,0xE8), False, False, BFONT)]], space_after=0)
def bubble(y, who, msg, h):
    if who == "bot":
        rect(s, pmx+0.25, y, 3.0, h, fill=CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        text(s, pmx+0.42, y+0.12, 2.7, h-0.2, [[(msg, 10.5, TXT, False, False, BFONT)]], line=1.08, space_after=0)
    else:
        rect(s, pmx+1.15, y, 3.0, h, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        text(s, pmx+1.32, y+0.12, 2.7, h-0.2, [[(msg, 10.5, WHITE, False, False, BFONT)]], line=1.08, space_after=0)
bubble(pmy+1.15, "bot", "Assalomu alaykum! Wall Street'ga xush kelibsiz \U0001F44B\nRo'yxatdan o'tamizmi?", 0.95)
bubble(pmy+2.25, "user", "Ha, ro'yxatdan o'taman ✅", 0.5)
bubble(pmy+2.90, "bot", "\U0001F4F1 Telefon raqamingizni yuboring:", 0.5)
bubble(pmy+3.55, "user", "+998 90 123 45 67", 0.5)
bubble(pmy+4.20, "bot", "Tayyor! Login: ali.valiyev\nParol: ••••••••\nIlovaga shu bilan kiring \U0001F389", 0.95)

# ============================================================
# SLIDE 15 — MOBIL ILOVA (kod mockup)
# ============================================================
s = slide(); bg(s, BG2); corner_dots(s)
num_badge(s, 0.7, 0.55, 15)
text(s, 1.55, 0.56, 9, 0.4, [[("MOBIL ILOVA — FLUTTER", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 7, 1.0, [[("Talaba uchun mobil ilova", 30, TXT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 2.15, 6.4, 1.0,
     [[("Talaba o'qish jarayonini telefonida kuzatadi. REST API (JWT) orqali tizim "
        "bilan xavfsiz bog'lanadi.", 13.5, MUT, False, False, BFONT)]], line=1.18)
mob_items = [
    ("Profil", "— shaxsiy ma'lumotlarni ko'rish."),
    ("Mening kurslarim", "— daraja, o'qituvchi, foiz."),
    ("To'lovlar", "— to'langan, qarz, xulosa."),
    ("Coin balansi", "— rag'batlantirish ballari."),
    ("Yangiliklar", "— markaz e'lonlari."),
]
bullets(s, 0.7, 3.4, 6.4, mob_items, gap=0.58, head_size=13.5, body_size=11.5)

def phone(px, py, header, accent, rows, balance=None):
    pw, ph = 2.7, 5.7
    rect(s, px, py, pw, ph, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, line=accent, line_w=1.5, shadow=True)
    rect(s, px, py, pw, 1.1, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    rect(s, px, py+0.55, pw, 0.55, fill=TEAL)
    text(s, px+0.25, py+0.34, pw-0.5, 0.45, [[(header, 13, WHITE, True, False, HFONT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    cy = py + 1.35
    if balance:
        rect(s, px+0.22, cy, pw-0.44, 1.0, fill=BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12, line=GOLD, line_w=1.0)
        text(s, px+0.22, cy+0.12, pw-0.44, 0.4, [[(balance[0], 11, MUT, False, False, BFONT)]], align=PP_ALIGN.CENTER, space_after=0)
        text(s, px+0.22, cy+0.4, pw-0.44, 0.55, [[(balance[1], 26, GOLD, True, False, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        cy += 1.2
    for title, sub, val in rows:
        rect(s, px+0.22, cy, pw-0.44, 0.78, fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12, line=TEAL, line_w=0.75)
        text(s, px+0.4, cy+0.12, pw-1.0, 0.3, [[(title, 11.5, TXT, True, False, HFONT)]], space_after=0)
        text(s, px+0.4, cy+0.42, pw-1.0, 0.3, [[(sub, 9.5, MUT, False, False, BFONT)]], space_after=0)
        if val:
            text(s, px+0.22, cy+0.22, pw-0.6, 0.35, [[(val, 12, accent, True, False, HFONT)]], align=PP_ALIGN.RIGHT, space_after=0)
        cy += 0.92

phone(7.55, 1.2, "Profil & Kurslarim", MINT,
      [("Ali Valiyev", "ID: WS-014 · Faol", ""),
       ("General English", "Pre-Intermediate", "82%"),
       ("IELTS 6.5", "Intensive", "67%"),
       ("Speaking Club", "Haftada 2 kun", "90%")])
phone(10.35, 1.2, "To'lov & Coin", GOLD,
      [("Iyun 2026", "General English", "OK"),
       ("May 2026", "To'langan", "OK"),
       ("Jami qarz", "Keyingi oy", "0")],
      balance=("Coin balansi", "1 240"))
text(s, 7.55, 7.02, 5.5, 0.3,
     [[("REST API:  /api/auth · /api/me · /api/me/courses · /api/me/payments · /api/news",
        9, MUT, False, True, BFONT)]], align=PP_ALIGN.CENTER, space_after=0)

# ============================================================
# SLIDE 16 — TEXNOLOGIYALAR VA NATIJALAR
# ============================================================
s = slide(); bg(s, BG); corner_dots(s)
num_badge(s, 0.7, 0.55, 16)
text(s, 1.55, 0.56, 9, 0.4, [[("TEXNOLOGIYALAR VA NATIJALAR", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 12, 0.7, [[("Qo'llanilgan texnologiyalar", 30, TXT, True, False, HFONT)]], space_after=0)
# tech stack rows
stacks = [
    ("Backend", "Python · Django · Django REST Framework · JWT", GOLD),
    ("Ma'lumotlar bazasi", "PostgreSQL (production) · SQLite (development)", MINT),
    ("Frontend", "Django Templates · Bootstrap · HTML/CSS · JavaScript", GOLD),
    ("Mobil", "Flutter · Dart · BLoC · JWT auth", MINT),
    ("Integratsiya", "aiogram (Telegram) · ReportLab (PDF) · QR · Excel eksport", GOLD),
    ("Deploy", "Render · PythonAnywhere · Git / GitHub", MINT),
]
ty = 2.1
for name, tools, acc in stacks:
    rect(s, 0.7, ty, 7.7, 0.66, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12, line=TEAL, line_w=0.75)
    rect(s, 0.7, ty, 0.1, 0.66, fill=acc, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, 1.0, ty, 2.4, 0.66, [[(name, 13, acc, True, False, HFONT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 3.3, ty, 5.0, 0.66, [[(tools, 11.5, TXT, False, False, BFONT)]], anchor=MSO_ANCHOR.MIDDLE, line=1.0, space_after=0)
    ty += 0.78
# stats column
stats = [("51", "talaba"), ("21", "kurs"), ("21", "o'qituvchi"), ("108", "to'lov yozuvi")]
sx = 8.75; sy = 2.1
text(s, sx, sy-0.02, 3.9, 0.4, [[("TIZIMDAGI MA'LUMOT", 11, MINT, True, False, HFONT)]], space_after=0)
sy += 0.5
for i, (num, lab) in enumerate(stats):
    r, c = divmod(i, 2)
    px = sx + c*2.0
    py = sy + r*1.55
    rect(s, px, py, 1.85, 1.4, fill=BG2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, line=TEAL, line_w=1.0)
    text(s, px, py+0.22, 1.85, 0.7, [[(num, 36, GOLD, True, False, HFONT)]], align=PP_ALIGN.CENTER, space_after=0)
    text(s, px, py+0.95, 1.85, 0.4, [[(lab, 12, MUT, False, False, BFONT)]], align=PP_ALIGN.CENTER, space_after=0)

# ============================================================
# SLIDE 17 — XULOSA
# ============================================================
s = slide(); bg(s, BG)
rect(s, -1.8, 4.0, 6.0, 6.0, fill=BG2, shape=MSO_SHAPE.OVAL)
rect(s, 10.5, -2.0, 5.5, 5.5, fill=CARD, shape=MSO_SHAPE.OVAL)
num_badge(s, 0.7, 0.55, 17)
text(s, 1.55, 0.56, 9, 0.4, [[("XULOSA", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 11.5, 0.8, [[("Ishning natijalari va xulosa", 30, TXT, True, False, HFONT)]], space_after=0)
concl = [
    ("Yagona platforma", "Talaba, kurs, davomat, to'lov va hisobot bitta tizimga birlashtirildi."),
    ("To'liq avtomatlashtirish", "To'lov, qarz va davomat foizi qo'lda hisobsiz, avtomatik aniqlanadi."),
    ("Rasmiy hujjatlar", "QR-kod va pechatli PDF cheklar hamda Excel/PDF hisobotlar."),
    ("Ko'p kanalli kirish", "Web CRM, mobil ilova va Telegram bot bitta baza bilan ishlaydi."),
    ("Kelajak rejasi", "SMS xabarnoma, onlayn to'lov shlyuzi va ko'p markazli rejim."),
]
cy = 2.25
for h, b in concl:
    rect(s, 0.85, cy+0.05, 0.34, 0.34, fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, 0.9, cy+0.02, 0.34, 0.34, [[("✓", 14, BG, True, False, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 1.4, cy, 10.8, 0.7,
         [[(h+" — ", 15, TXT, True, False, HFONT), (b, 13, MUT, False, False, BFONT)]],
         line=1.05, space_after=0)
    cy += 0.78
# thanks bar
rect(s, 0.7, 6.45, 11.93, 0.75, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2, line=GOLD, line_w=1.0)
text(s, 0.7, 6.45, 11.93, 0.75, [[("E'tiboringiz uchun rahmat!", 18, WHITE, True, False, HFONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))

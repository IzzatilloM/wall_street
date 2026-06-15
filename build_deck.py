# -*- coding: utf-8 -*-
"""Wall Street CRM — BMI himoya taqdimoti (20 slayd) generatori (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

SHOTS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_pptx_shots")
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Wall_Street_CRM_BMI.pptx")

# ── Palitra ───────────────────────────────────────────────
INK      = "0E2A2C"
DEEP     = "0A3A3D"
DEEPER   = "072A2C"
BRAND    = "0E6E66"
BRAND_DK = "0A4F49"
TEAL     = "17A398"
TEAL_TINT= "E4F1EF"
GOLD     = "E0A82F"
GOLD_TINT= "FAEFD3"
GREEN    = "1F9D67"
RED      = "D64545"
PURPLE   = "6D4FC3"
PURPLE_T = "ECE7FA"
LIGHT    = "F4F8F7"
WHITE    = "FFFFFF"
MUTED    = "5E7373"
MUTED2   = "8AA0A0"
LINE     = "DBE7E4"

HEAD = "Trebuchet MS"
BODY = "Calibri"

EMUIN = 914400
PAGEW, PAGEH = 13.333, 7.5
AR = 2880 / 1800  # 1.6

prs = Presentation()
prs.slide_width = Emu(int(PAGEW * EMUIN))
prs.slide_height = Emu(int(PAGEH * EMUIN))
BLANK = prs.slide_layouts[6]


def C(h):
    return RGBColor.from_string(h)


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(bg)
    return s


def set_shadow(shape, blur=11, dist=4, dir_deg=90, color="06312F", alpha=26):
    sp = shape._element.spPr
    for e in sp.findall(qn("a:effectLst")):
        sp.remove(e)
    eff = sp.makeelement(qn("a:effectLst"), {})
    sh = sp.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur * 12700)),
        "dist": str(int(dist * 12700)),
        "dir": str(int(dir_deg * 60000)),
        "rotWithShape": "0",
    })
    clr = sp.makeelement(qn("a:srgbClr"), {"val": color})
    al = sp.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
    clr.append(al)
    sh.append(clr)
    eff.append(sh)
    sp.append(eff)


def rrect(s, x, y, w, h, fill=WHITE, line=None, line_w=1.0, radius=0.12,
          shadow=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = C(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = C(line)
        sh.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
        except Exception:
            pass
    if shadow:
        set_shadow(sh)
    return sh


def oval(s, x, y, w, h, fill, line=None, line_w=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = C(line)
        sh.line.width = Pt(line_w)
    return sh


def line(s, x1, y1, x2, y2, color=TEAL, w=1.5, dash=None):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                Inches(x2), Inches(y2))
    cn.line.color.rgb = C(color)
    cn.line.width = Pt(w)
    cn.shadow.inherit = False
    if dash:
        ln = cn.line._get_or_add_ln()
        pd = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(pd)
    return cn


def txt(s, runs, x, y, w, h, size=14, color=INK, bold=False, font=BODY,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, wrap=True, spacing=1.0,
        space_after=0.0, margin=0.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(margin))
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        if space_after:
            p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for t, o in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", bold)
            r.font.italic = o.get("italic", False)
            r.font.name = o.get("font", font)
            r.font.color.rgb = C(o.get("color", color))
            if "spc" in o:
                r.font._rPr.set("spc", str(int(o["spc"] * 100)))
    return tb


def chip(s, x, y, size, num, fill=BRAND, tcolor=WHITE):
    rrect(s, x, y, size, size, fill=fill, radius=0.14, shadow=True)
    txt(s, num, x, y - 0.02, size, size, size=21, color=tcolor, bold=True,
        font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def header(s, num, title, kicker, color=BRAND, dark=False):
    tcol = WHITE if dark else INK
    kcol = color if not dark else TEAL
    chip(s, 0.55, 0.5, 0.66, num, fill=color)
    txt(s, kicker, 1.42, 0.46, 11.0, 0.3, size=11.5, color=kcol, bold=True,
        font=HEAD, spacing=1.0, margin=0.0)
    txt(s, [(title, {})], 1.40, 0.7, 11.4, 0.85, size=29, color=tcol, bold=True,
        font=HEAD, valign=MSO_ANCHOR.TOP, margin=0.0)


def footer(s, idx, dark=False):
    col = MUTED2 if not dark else "5F7E7E"
    # logo glyph
    rrect(s, 0.55, 7.04, 0.26, 0.26, fill=(BRAND if not dark else TEAL), radius=0.07)
    txt(s, "W", 0.55, 7.0, 0.26, 0.26, size=12, color=GOLD, bold=True,
        font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    txt(s, [[("WALL STREET", {"bold": True, "color": col}), ("  ·  TA'LIM CRM", {"color": col})]],
        0.9, 7.04, 5.0, 0.28, size=8.5, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    txt(s, f"{idx:02d} / 20", 10.8, 7.04, 1.98, 0.28, size=9, color=col,
        bold=True, font=HEAD, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, margin=0.0)


def img_card(s, fname, x, y, w, h, caption=None, pad=0.13, accent=BRAND):
    rrect(s, x, y, w, h, fill=WHITE, radius=0.12, shadow=True)
    cap_h = 0.42 if caption else 0.0
    iw_max = w - 2 * pad
    ih_max = h - 2 * pad - cap_h
    if iw_max / ih_max > AR:
        ih = ih_max
        iw = ih * AR
    else:
        iw = iw_max
        ih = iw / AR
    ix = x + (w - iw) / 2
    iy = y + pad + (ih_max - ih) / 2
    rrect(s, ix - 0.02, iy - 0.02, iw + 0.04, ih + 0.04, fill=None,
          line=LINE, line_w=1.0, radius=0.05)
    s.shapes.add_picture(os.path.join(SHOTS, fname), Inches(ix), Inches(iy),
                         Inches(iw), Inches(ih))
    if caption:
        txt(s, caption, x + pad, y + h - cap_h, w - 2 * pad, cap_h, size=11,
            color=MUTED, bold=True, font=HEAD, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE)


def logo_mark(s, x, y, size, dark=False):
    rrect(s, x, y, size, size, fill=(WHITE if dark else BRAND), radius=size * 0.22,
          shadow=True)
    txt(s, "W", x, y - 0.02, size, size, size=int(size * 34), color=GOLD,
        bold=True, font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
# SLAYD 1 — Sarlavha
# ════════════════════════════════════════════════════════════
def slide1():
    s = add_slide(DEEPER)
    # dekorativ panellar
    rrect(s, 8.7, -1.4, 6.5, 6.5, fill=DEEP, radius=1.2)
    oval(s, 10.7, 4.6, 4.4, 4.4, fill=BRAND_DK)
    rrect(s, -1.0, 5.6, 6.0, 3.0, fill=DEEP, radius=0.8)
    # logo + wordmark
    logo_mark(s, 0.95, 0.85, 0.95, dark=True)
    txt(s, [[("Wall Street", {"bold": True, "color": WHITE})]], 2.05, 0.86, 6, 0.5,
        size=23, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    txt(s, [[("TA'LIM CRM", {"bold": True, "color": TEAL, "spc": 2.5})]], 2.07, 1.34, 6, 0.32,
        size=11, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    # kicker
    rrect(s, 0.97, 2.7, 6.2, 0.46, fill=BRAND, radius=0.23)
    txt(s, [[("BITIRUV MALAKAVIY ISHI  ·  HIMOYA TAQDIMOTI", {"spc": 1.5})]],
        1.2, 2.7, 6.0, 0.46, size=11.5, color=WHITE, bold=True, font=HEAD,
        valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    # asosiy sarlavha
    txt(s, [[("Wall Street CRM", {})]], 0.92, 3.25, 11.5, 1.2, size=54, color=WHITE,
        bold=True, font=HEAD, margin=0.0)
    txt(s, [[("Ta'lim markazlari uchun ", {"color": "CFE6E2"}),
             ("sun'iy intellektli", {"color": GOLD, "bold": True}),
             (" yagona boshqaruv tizimi", {"color": "CFE6E2"})]],
        0.95, 4.45, 10.8, 0.6, size=19, font=BODY, margin=0.0)
    # meta blok
    rrect(s, 0.95, 5.5, 11.4, 1.25, fill=DEEP, radius=0.14)
    metas = [("Bajardi", "________________"), ("Ilmiy rahbar", "________________"),
             ("Yo'nalish", "________________"), ("Yil", "Toshkent — 2026")]
    mw = 11.4 / 4
    for i, (k, v) in enumerate(metas):
        mx = 0.95 + i * mw
        txt(s, k.upper(), mx + 0.3, 5.66, mw - 0.4, 0.3, size=10, color=TEAL,
            bold=True, font=HEAD, margin=0.0)
        txt(s, v, mx + 0.3, 5.98, mw - 0.4, 0.55, size=14, color=WHITE,
            bold=True, font=BODY, margin=0.0)
        if i < 3:
            line(s, mx + mw, 5.72, mx + mw, 6.5, color=BRAND_DK, w=1.0)


# ════════════════════════════════════════════════════════════
# SLAYD 2 — Mundarija
# ════════════════════════════════════════════════════════════
def slide2():
    s = add_slide(LIGHT)
    header(s, "•", "Mundarija", "TAQDIMOT REJASI")
    items = [
        "Muammo va dolzarblik", "Loyiha maqsadi va vazifalari",
        "Tizim arxitekturasi va texnologiyalar", "Platforma imkoniyatlari",
        "Boshqaruv paneli (Dashboard)", "Talabalar boshqaruvi",
        "AI — talaba churn (ketib qolish) bashorati", "Kurslar va talaba transferi",
        "Ro'yxatga olish (Enrollments)", "AI — kurs tavsiyasi",
        "O'qituvchilar boshqaruvi", "AI — o'qituvchi samaradorlik hisoboti",
        "AI ustunligi — boshqa CRM bilan solishtirish", "Davomat, to'lov, coin va reyting",
        "Telegram botlar (2 ta)", "O'qituvchi mobil ilovasi",
        "Talaba mobil ilovasi", "Xulosa va natijalar",
    ]
    col_w = 5.85
    x0 = [0.7, 6.95]
    y0 = 1.85
    rh = 0.555
    for i, it in enumerate(items):
        col = 0 if i < 9 else 1
        row = i if i < 9 else i - 9
        x = x0[col]
        y = y0 + row * rh
        rrect(s, x, y, col_w, 0.5, fill=WHITE, radius=0.1, shadow=True)
        rrect(s, x + 0.09, y + 0.09, 0.32, 0.32, fill=(PURPLE if "AI" in it else BRAND),
              radius=0.07)
        txt(s, f"{i+1:02d}", x + 0.09, y + 0.07, 0.32, 0.32, size=10.5, color=WHITE,
            bold=True, font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, it, x + 0.55, y, col_w - 0.65, 0.5, size=12.5, color=INK,
            bold=True, font=BODY, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    footer(s, 2)


# ════════════════════════════════════════════════════════════
# SLAYD 3 — Muammo
# ════════════════════════════════════════════════════════════
def slide3():
    s = add_slide(LIGHT)
    header(s, "01", "Muammo va dolzarblik", "NEGA YANGI YECHIM KERAK?")
    txt(s, [[("An'anaviy o'quv markazlari ma'lumotlarni daftar, Excel va tarqoq "
              "ilovalarda yuritadi — bu vaqt yo'qotish, xato va nazoratsizlikka olib keladi.",
              {})]], 0.7, 1.62, 11.9, 0.6, size=14.5, color=MUTED, font=BODY, margin=0.0)
    probs = [
        ("01", "Qo'lda yuritish", "Davomat va to'lovlar daftar/Excelda — sekin va xatoga moyil.", RED),
        ("02", "Ko'rinmas churn", "Talabaning ketib qolishi oldindan sezilmaydi, kech qoladi.", GOLD),
        ("03", "Subyektiv baho", "O'qituvchi samaradorligi raqamlarsiz, his bilan baholanadi.", PURPLE),
        ("04", "Tarqoq ma'lumot", "Telegram, qog'oz, jadval — yagona manba yo'q.", BRAND),
    ]
    y = 2.35
    for i, (n, t, d, col) in enumerate(probs):
        yy = y + i * 1.07
        rrect(s, 0.7, yy, 11.9, 0.92, fill=WHITE, radius=0.12, shadow=True)
        rrect(s, 0.7, yy, 0.13, 0.92, fill=col, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
        oval(s, 1.05, yy + 0.21, 0.5, 0.5, fill=col)
        txt(s, n, 1.05, yy + 0.19, 0.5, 0.5, size=15, color=WHITE, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, t, 1.8, yy + 0.13, 3.4, 0.66, size=16, color=INK, bold=True,
            font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, d, 5.3, yy + 0.13, 7.1, 0.66, size=13, color=MUTED, font=BODY,
            valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    footer(s, 3)


# ════════════════════════════════════════════════════════════
# SLAYD 4 — Maqsad va vazifalar
# ════════════════════════════════════════════════════════════
def slide4():
    s = add_slide(LIGHT)
    header(s, "02", "Loyiha maqsadi va vazifalari", "YECHIM")
    rrect(s, 0.7, 1.65, 11.9, 1.0, fill=BRAND, radius=0.14, shadow=True)
    txt(s, [[("MAQSAD   ", {"color": GOLD, "bold": True, "size": 12, "font": HEAD}),
             ("Talabalar, kurslar, to'lovlar va davomatni AI bilan kuchaytirilgan "
              "yagona veb hamda mobil platformada boshqarish.", {"color": WHITE, "size": 15})]],
        1.0, 1.65, 11.3, 1.0, valign=MSO_ANCHOR.MIDDLE, margin=0.0, spacing=1.05)
    tasks = [
        ("Yagona boshqaruv", "Talaba, kurs, guruh va ro'yxatga olishni bitta tizimda yuritish."),
        ("Moliya nazorati", "To'lovlar, qarzlar va o'qituvchi maoshini avtomatik hisoblash."),
        ("Avtomat davomat", "Davomat olish, hisobot va real vaqt statistikasi."),
        ("AI tahlil", "Churn bashorati, kurs tavsiyasi va o'qituvchi hisoboti."),
        ("2 Telegram bot", "Ro'yxatdan o'tkazish va texnik yordam botlari."),
        ("Mobil ilovalar", "O'qituvchi va talaba uchun alohida Flutter ilovalari."),
    ]
    cw, ch = 3.82, 1.32
    gx, gy = 0.24, 0.26
    x0, y0 = 0.7, 2.95
    for i, (t, d) in enumerate(tasks):
        c = i % 3
        r = i // 3
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        rrect(s, x, y, cw, ch, fill=WHITE, radius=0.12, shadow=True)
        oval(s, x + 0.22, y + 0.24, 0.42, 0.42, fill=TEAL_TINT)
        txt(s, "✓", x + 0.22, y + 0.21, 0.42, 0.42, size=16, color=BRAND, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, t, x + 0.78, y + 0.2, cw - 0.95, 0.4, size=15, color=INK, bold=True,
            font=HEAD, margin=0.0)
        txt(s, d, x + 0.26, y + 0.66, cw - 0.5, 0.6, size=12, color=MUTED, font=BODY,
            margin=0.0, spacing=1.0)
    footer(s, 4)


# ════════════════════════════════════════════════════════════
# SLAYD 5 — Arxitektura
# ════════════════════════════════════════════════════════════
def slide5():
    s = add_slide(LIGHT)
    header(s, "03", "Tizim arxitekturasi va texnologiyalar", "QANDAY QURILGAN")
    cx, cy, cw, chh = 5.27, 3.15, 2.8, 1.35
    ccx, ccy = cx + cw / 2, cy + chh / 2
    sats = [
        (1.0, 1.75, 3.2, 1.05, "Ma'lumotlar bazasi", "SQLite / PostgreSQL", BRAND),
        (9.1, 1.75, 3.25, 1.05, "Claude AI", "Anthropic — tahlil va bashorat", PURPLE),
        (1.0, 5.0, 3.35, 1.1, "Mobil ilovalar", "O'qituvchi · Talaba (Flutter)", TEAL),
        (8.95, 5.0, 3.4, 1.1, "Telegram botlar", "Education · Technic (aiogram)", GOLD),
    ]
    for x, y, w, h, t, d, col in sats:
        bx, by = x + w / 2, y + h / 2
        line(s, ccx, ccy, bx, by, color="B9CFCB", w=1.6)
    for x, y, w, h, t, d, col in sats:
        rrect(s, x, y, w, h, fill=WHITE, radius=0.12, shadow=True, line=col, line_w=1.5)
        rrect(s, x, y, 0.12, h, fill=col, shape=MSO_SHAPE.RECTANGLE)
        txt(s, t, x + 0.3, y + 0.16, w - 0.45, 0.4, size=15, color=INK, bold=True,
            font=HEAD, margin=0.0)
        txt(s, d, x + 0.3, y + 0.56, w - 0.45, 0.45, size=11.5, color=MUTED, font=BODY,
            margin=0.0)
    # markaz
    rrect(s, cx, cy, cw, chh, fill=BRAND, radius=0.14, shadow=True)
    txt(s, [[("Wall Street CRM", {"color": WHITE, "bold": True, "size": 16})],
            [("Django · REST API · JWT", {"color": "BFE0DC", "size": 11.5})]],
        cx, cy, cw, chh, font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        margin=0.0, spacing=1.05)
    # texnologiya pill
    stack = "Python   ·   Django   ·   DRF   ·   JWT   ·   PostgreSQL   ·   aiogram   ·   Anthropic Claude   ·   Flutter"
    rrect(s, 1.4, 6.45, 10.55, 0.5, fill=INK, radius=0.25)
    txt(s, stack, 1.4, 6.45, 10.55, 0.5, size=11.5, color=WHITE, bold=True,
        font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    footer(s, 5)


# ════════════════════════════════════════════════════════════
# SLAYD 6 — Platforma imkoniyatlari
# ════════════════════════════════════════════════════════════
def slide6():
    s = add_slide(LIGHT)
    header(s, "04", "Platforma imkoniyatlari", "YAGONA TIZIMDA 13+ MODUL")
    img_card(s, "01_login.png", 0.7, 1.7, 5.55, 3.55, caption="Xavfsiz kirish · ko'p tilli interfeys")
    mods = ["Boshqaruv paneli", "Talabalar", "Kurslar", "Ro'yxatga olish",
            "O'qituvchilar", "Davomat", "To'lovlar", "Maosh",
            "Hisobotlar", "Kalendar", "Sozlamalar", "Yangiliklar"]
    x0, y0 = 6.5, 1.7
    cw, ch = 1.93, 0.62
    gx, gy = 0.13, 0.13
    for i, m in enumerate(mods):
        c = i % 3
        r = i // 3
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        rrect(s, x, y, cw, ch, fill=WHITE, radius=0.1, shadow=True)
        oval(s, x + 0.13, y + 0.16, 0.3, 0.3, fill=TEAL_TINT)
        oval(s, x + 0.205, y + 0.235, 0.15, 0.15, fill=BRAND)
        txt(s, m, x + 0.52, y, cw - 0.6, ch, size=11, color=INK, bold=True,
            font=BODY, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    # mini thumbnaillar
    minis = [("14_salary.png", "Maosh"), ("15_calendar.png", "Kalendar"),
             ("17_settings.png", "Sozlamalar")]
    mx0 = 6.5
    mw = 1.93
    for i, (f, cap) in enumerate(minis):
        x = mx0 + i * (mw + 0.15)
        img_card(s, f, x, 4.78, mw, 1.62, caption=cap, pad=0.08)
    footer(s, 6)


# ════════════════════════════════════════════════════════════
# Umumiy: katta skrinshot + funksiya paneli
# ════════════════════════════════════════════════════════════
def feature_slide(idx, num, title, kicker, shot, feats, shotcap, color=BRAND,
                  second=None, advantage=None):
    s = add_slide(LIGHT)
    header(s, num, title, kicker, color=color)
    big_w = 7.35 if not second else 7.35
    img_card(s, shot, 0.7, 1.72, big_w, 4.95 if not second else 3.35, caption=shotcap,
             accent=color)
    if second:
        img_card(s, second[0], 0.7, 5.2, big_w, 1.47, caption=second[1], pad=0.08,
                 accent=color)
    px = 0.7 + big_w + 0.3
    pw = PAGEW - px - 0.55
    ty = 1.72
    txt(s, "IMKONIYATLAR", px, ty, pw, 0.3, size=11, color=color, bold=True, font=HEAD,
        margin=0.0)
    yy = ty + 0.42
    for t, d in feats:
        rrect(s, px, yy, pw, 0.86, fill=WHITE, radius=0.1, shadow=True)
        oval(s, px + 0.18, yy + 0.27, 0.34, 0.34, fill=color)
        txt(s, "✓", px + 0.18, yy + 0.25, 0.34, 0.34, size=13, color=WHITE, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, t, px + 0.64, yy + 0.1, pw - 0.78, 0.34, size=13, color=INK, bold=True,
            font=HEAD, margin=0.0)
        txt(s, d, px + 0.64, yy + 0.44, pw - 0.78, 0.36, size=10.5, color=MUTED, font=BODY,
            margin=0.0)
        yy += 0.98
    if advantage:
        rrect(s, px, yy + 0.02, pw, 6.65 - yy, fill=color, radius=0.12, shadow=True)
        txt(s, [[("AFZALLIK   ", {"color": GOLD, "bold": True, "size": 10.5, "font": HEAD})],
                [(advantage, {"color": WHITE, "size": 11.5})]],
            px + 0.22, yy + 0.12, pw - 0.4, 6.65 - yy - 0.2, font=BODY,
            valign=MSO_ANCHOR.MIDDLE, margin=0.0, spacing=1.05)
    footer(s, idx)
    return s


# ════════════════════════════════════════════════════════════
# AI slayd shabloni (binafsha urg'u + afzallik banneri)
# ════════════════════════════════════════════════════════════
def ai_slide(idx, num, title, shot, how, advantage, shotcap):
    s = add_slide(LIGHT)
    header(s, num, title, "SUN'IY INTELLEKT", color=PURPLE)
    # AI rozet
    rrect(s, 10.95, 0.52, 1.85, 0.56, fill=PURPLE, radius=0.28, shadow=True)
    txt(s, [[("🤖  Claude AI", {})]], 10.95, 0.52, 1.85, 0.56, size=12, color=WHITE,
        bold=True, font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    img_card(s, shot, 0.7, 1.72, 7.5, 4.95, caption=shotcap, accent=PURPLE)
    px = 8.5
    pw = PAGEW - px - 0.55
    txt(s, "QANDAY ISHLAYDI", px, 1.72, pw, 0.3, size=11, color=PURPLE, bold=True,
        font=HEAD, margin=0.0)
    yy = 2.12
    for i, (t, d) in enumerate(how):
        rrect(s, px, yy, pw, 1.0, fill=WHITE, radius=0.1, shadow=True)
        rrect(s, px + 0.16, yy + 0.16, 0.36, 0.36, fill=PURPLE_T, radius=0.08)
        txt(s, str(i + 1), px + 0.16, yy + 0.14, 0.36, 0.36, size=13, color=PURPLE,
            bold=True, font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, t, px + 0.64, yy + 0.13, pw - 0.78, 0.34, size=12.5, color=INK, bold=True,
            font=HEAD, margin=0.0)
        txt(s, d, px + 0.64, yy + 0.46, pw - 0.78, 0.46, size=10.5, color=MUTED, font=BODY,
            margin=0.0, spacing=1.0)
        yy += 1.1
    rrect(s, px, yy + 0.02, pw, 6.67 - yy, fill=PURPLE, radius=0.12, shadow=True)
    txt(s, [[("✦  BOSHQA CRM'LARDA YO'Q", {"color": GOLD, "bold": True, "size": 10.5, "font": HEAD})],
            [(advantage, {"color": WHITE, "size": 11})]],
        px + 0.22, yy + 0.1, pw - 0.42, 6.67 - yy - 0.16, font=BODY,
        valign=MSO_ANCHOR.MIDDLE, margin=0.0, spacing=1.02)
    footer(s, idx)


# ════════════════════════════════════════════════════════════
# SLAYD 15 — AI solishtirish (qora slayd)
# ════════════════════════════════════════════════════════════
def slide15():
    s = add_slide(DEEPER)
    rrect(s, 9.3, -1.2, 6.0, 5.5, fill=DEEP, radius=1.0)
    header(s, "13", "AI ustunligi — boshqa CRM bilan solishtirish", "NEGA WALL STREET CRM?",
           color=TEAL, dark=True)
    rows = [
        ("AI churn (ketib qolish) bashorati", True, False),
        ("AI kurs tavsiyasi", True, False),
        ("AI o'qituvchi samaradorlik hisoboti", True, False),
        ("Yuqori-risk avtomatik ogohlantirish", True, False),
        ("Coin & reyting (gamifikatsiya)", True, "qisman"),
        ("2 ta Telegram bot integratsiyasi", True, "kam"),
        ("O'qituvchi + Talaba mobil ilovasi", True, "qisman"),
        ("Talabani kursdan-kursga transfer", True, False),
    ]
    x1, w1 = 0.7, 7.05
    x2, w2 = 7.95, 2.35
    x3, w3 = 10.45, 2.35
    yh = 1.78
    # ustun sarlavhalari
    txt(s, "IMKONIYAT", x1 + 0.2, yh, w1, 0.4, size=11, color=MUTED2, bold=True, font=HEAD,
        valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    rrect(s, x2, yh, w2, 0.5, fill=BRAND, radius=0.1)
    txt(s, "Wall Street CRM", x2, yh, w2, 0.5, size=12, color=WHITE, bold=True, font=HEAD,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    rrect(s, x3, yh, w3, 0.5, fill="3A4F50", radius=0.1)
    txt(s, "An'anaviy CRM", x3, yh, w3, 0.5, size=12, color="C8D6D5", bold=True, font=HEAD,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    y = yh + 0.62
    rh = 0.585
    for i, (label, ws, other) in enumerate(rows):
        yy = y + i * rh
        rrect(s, x1, yy, w1, 0.5, fill=("0F4548" if i % 2 == 0 else "0C3C3F"), radius=0.08)
        txt(s, label, x1 + 0.2, yy, w1 - 0.3, 0.5, size=12.5, color="EAF3F2", bold=True,
            font=BODY, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        # WS ustun — ✓
        rrect(s, x2, yy, w2, 0.5, fill="123E3C", radius=0.08, line=BRAND, line_w=1.0)
        oval(s, x2 + w2 / 2 - 0.61, yy + 0.1, 0.3, 0.3, fill=GREEN)
        txt(s, "✓", x2 + w2 / 2 - 0.61, yy + 0.08, 0.3, 0.3, size=12, color=WHITE, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, "Bor", x2 + w2 / 2 - 0.22, yy, 0.85, 0.5, size=12, color=WHITE, bold=True,
            font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        # boshqa ustun
        rrect(s, x3, yy, w3, 0.5, fill="20302F", radius=0.08)
        if other is False:
            oval(s, x3 + w3 / 2 - 0.58, yy + 0.1, 0.3, 0.3, fill=RED)
            txt(s, "✕", x3 + w3 / 2 - 0.58, yy + 0.08, 0.3, 0.3, size=12, color=WHITE,
                bold=True, font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
            txt(s, "Yo'q", x3 + w3 / 2 - 0.19, yy, 0.9, 0.5, size=12, color="9DB0AF",
                bold=True, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        else:
            txt(s, other, x3, yy, w3, 0.5, size=11.5, color=GOLD, bold=True, font=HEAD,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    footer(s, 15, dark=True)


# ════════════════════════════════════════════════════════════
# SLAYD 16 — Operatsion modullar (2x2)
# ════════════════════════════════════════════════════════════
def slide16():
    s = add_slide(LIGHT)
    header(s, "14", "Davomat, to'lov, coin va reyting", "KUNDALIK OPERATSIYALAR")
    grid = [("08_attendance.png", "Davomat olish va statistika", 0.7, 1.72),
            ("10_payments.png", "To'lovlar va qarzlar nazorati", 6.75, 1.72),
            ("09_coin_award.png", "Coin berish · gamifikatsiya", 0.7, 4.3),
            ("13_reports.png", "Hisobotlar va tahlil", 6.75, 4.3)]
    for f, cap, x, y in grid:
        img_card(s, f, x, y, 5.88, 2.46, caption=cap, pad=0.1)
    footer(s, 16)


# ════════════════════════════════════════════════════════════
# SLAYD 17 — Telegram botlar
# ════════════════════════════════════════════════════════════
def slide17():
    s = add_slide(LIGHT)
    header(s, "15", "Telegram botlar (2 ta)", "AVTOMATLASHTIRISH", color=GOLD)
    # Bot 1
    rrect(s, 0.7, 1.75, 5.85, 3.05, fill=WHITE, radius=0.14, shadow=True)
    rrect(s, 0.7, 1.75, 5.85, 0.78, fill=BRAND, radius=0.14)
    rrect(s, 0.7, 2.2, 5.85, 0.33, fill=BRAND, shape=MSO_SHAPE.RECTANGLE)
    txt(s, [[("①  wall_street_education_bot", {})]], 1.0, 1.75, 5.4, 0.78, size=15,
        color=WHITE, bold=True, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    steps = [("Ro'yxatdan o'tish", "Foydalanuvchi bot orqali ariza qoldiradi."),
             ("Tizimga tushadi", "Ma'lumot avtomatik CRM bazasiga keladi."),
             ("Admin tasdiqlaydi", "Admin ko'rib chiqib, tasdiqlaydi yoki rad etadi."),
             ("Bildirishnoma", "Natija foydalanuvchiga Telegramga yuboriladi.")]
    yy = 2.66
    for i, (t, d) in enumerate(steps):
        rrect(s, 0.95, yy, 0.32, 0.32, fill=BRAND, radius=0.07)
        txt(s, str(i + 1), 0.95, yy - 0.01, 0.32, 0.32, size=11, color=WHITE, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, [[(t + " — ", {"bold": True, "color": INK, "size": 12.5}),
                 (d, {"color": MUTED, "size": 11.5})]],
            1.4, yy - 0.04, 4.95, 0.45, font=BODY, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        yy += 0.5
    # Bot 2
    rrect(s, 6.75, 1.75, 5.85, 3.05, fill=WHITE, radius=0.14, shadow=True)
    rrect(s, 6.75, 1.75, 5.85, 0.78, fill=GOLD, radius=0.14)
    rrect(s, 6.75, 2.2, 5.85, 0.33, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    txt(s, [[("②  wall_street_technic_bot", {})]], 7.05, 1.75, 5.4, 0.78, size=15,
        color=WHITE, bold=True, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    steps2 = [("Texnik yordam", "Foydalanuvchi muammo/murojaat yuboradi."),
              ("Ticket yaratiladi", "Murojaat «Yordam markazi» bo'limiga tushadi."),
              ("Admin javobi", "Admin javob yozadi, holatni boshqaradi."),
              ("Tezkor aloqa", "Javob foydalanuvchiga Telegramga qaytadi.")]
    yy = 2.66
    for i, (t, d) in enumerate(steps2):
        rrect(s, 7.0, yy, 0.32, 0.32, fill=GOLD, radius=0.07)
        txt(s, str(i + 1), 7.0, yy - 0.01, 0.32, 0.32, size=11, color=WHITE, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, [[(t + " — ", {"bold": True, "color": INK, "size": 12.5}),
                 (d, {"color": MUTED, "size": 11.5})]],
            7.45, yy - 0.04, 4.95, 0.45, font=BODY, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        yy += 0.5
    # pastki: registratsiya skrinshoti + afzallik
    img_card(s, "02_register.png", 0.7, 5.0, 5.85, 1.67, caption="Ro'yxatdan o'tish sahifasi",
             pad=0.08)
    rrect(s, 6.75, 5.0, 5.85, 1.67, fill=INK, radius=0.14, shadow=True)
    txt(s, [[("AFZALLIK", {"color": GOLD, "bold": True, "size": 11, "font": HEAD})],
            [("Botlar orqali ro'yxatga olish va texnik yordam to'liq avtomatlashgan — "
              "barcha murojaatlar bitta tizimda yig'iladi, admin Telegramni tark etmasdan "
              "boshqaradi.", {"color": "DCEDEB", "size": 12.5})]],
        7.05, 5.0, 5.3, 1.67, font=BODY, valign=MSO_ANCHOR.MIDDLE, margin=0.0, spacing=1.08)
    footer(s, 17)


# ════════════════════════════════════════════════════════════
# Mobil ilova slaydlari (bo'sh telefon ramkalari)
# ════════════════════════════════════════════════════════════
def phone_frame(s, x, y, w, h, label, color):
    rrect(s, x, y, w, h, fill=WHITE, radius=0.22, shadow=True, line=color, line_w=2.0)
    # ekran
    rrect(s, x + 0.12, y + 0.16, w - 0.24, h - 0.32, fill=LIGHT, radius=0.16,
          line=LINE, line_w=1.0)
    # notch
    rrect(s, x + w / 2 - 0.32, y + 0.12, 0.64, 0.12, fill=color, radius=0.06)
    # kamera placeholder belgisi
    ic = y + (h - 0.32) / 2 - 0.1
    oval(s, x + w / 2 - 0.28, ic - 0.2, 0.56, 0.56, fill="FFFFFF", line=color, line_w=1.5)
    oval(s, x + w / 2 - 0.13, ic - 0.05, 0.26, 0.26, fill=color)
    txt(s, [[("Skrinshot", {"bold": True, "color": MUTED, "size": 11.5})],
            [("shu yerga qo'yiladi", {"color": MUTED2, "size": 10})]],
        x + 0.1, ic + 0.5, w - 0.2, 0.7, font=BODY, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.TOP, margin=0.0, spacing=1.0)
    # caption pill
    rrect(s, x + 0.2, y + h - 0.62, w - 0.4, 0.42, fill=color, radius=0.21)
    txt(s, label, x + 0.2, y + h - 0.62, w - 0.4, 0.42, size=11, color=WHITE, bold=True,
        font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)


def mobile_slide(idx, num, title, subtitle, color, feats, frames):
    s = add_slide(LIGHT)
    header(s, num, title, "FLUTTER MOBIL ILOVA", color=color)
    rrect(s, 10.7, 0.5, 2.1, 0.58, fill=color, radius=0.29, shadow=True)
    txt(s, "📱  Flutter", 10.7, 0.5, 2.1, 0.58, size=12, color=WHITE, bold=True, font=HEAD,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    txt(s, subtitle, 0.7, 1.62, 11.9, 0.45, size=14, color=MUTED, font=BODY, margin=0.0)
    # feature chiplar (chap)
    fx, fw = 0.7, 3.55
    txt(s, "IMKONIYATLAR", fx, 2.2, fw, 0.3, size=11, color=color, bold=True, font=HEAD,
        margin=0.0)
    yy = 2.6
    for t, d in feats:
        rrect(s, fx, yy, fw, 0.92, fill=WHITE, radius=0.1, shadow=True)
        oval(s, fx + 0.18, yy + 0.29, 0.34, 0.34, fill=color)
        txt(s, "✓", fx + 0.18, yy + 0.27, 0.34, 0.34, size=13, color=WHITE, bold=True,
            font=HEAD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, t, fx + 0.64, yy + 0.13, fw - 0.78, 0.34, size=13, color=INK, bold=True,
            font=HEAD, margin=0.0)
        txt(s, d, fx + 0.64, yy + 0.46, fw - 0.78, 0.38, size=10.5, color=MUTED, font=BODY,
            margin=0.0)
        yy += 1.02
    # telefon ramkalari (o'ng)
    n = len(frames)
    area_x, area_w = 4.55, PAGEW - 4.55 - 0.55
    gap = 0.35
    pw = min(2.1, (area_w - (n - 1) * gap) / n)
    total = n * pw + (n - 1) * gap
    startx = area_x + (area_w - total) / 2
    py, ph = 2.2, 4.4
    for i, lab in enumerate(frames):
        x = startx + i * (pw + gap)
        phone_frame(s, x, py, pw, ph, lab, color)
    footer(s, idx)


# ════════════════════════════════════════════════════════════
# SLAYD 20 — Xulosa
# ════════════════════════════════════════════════════════════
def slide20():
    s = add_slide(DEEPER)
    rrect(s, 8.9, -1.3, 6.5, 6.0, fill=DEEP, radius=1.1)
    oval(s, 11.0, 4.8, 4.0, 4.0, fill=BRAND_DK)
    logo_mark(s, 0.7, 0.6, 0.8, dark=True)
    txt(s, [[("XULOSA VA NATIJALAR", {"spc": 2.0})]], 1.7, 0.62, 8, 0.8, size=13,
        color=TEAL, bold=True, font=HEAD, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    txt(s, [[("Wall Street CRM — to'liq, ", {"color": WHITE}),
             ("AI bilan kuchaytirilgan", {"color": GOLD, "bold": True}),
             (" ta'lim platformasi", {"color": WHITE})]],
        0.7, 1.55, 11.5, 1.0, size=30, bold=True, font=HEAD, margin=0.0)
    # natija statlar
    stats = [("3", "ta noyob AI funksiyasi"), ("2", "Telegram bot"),
             ("2", "mobil ilova (o'qituvchi · talaba)"), ("13+", "boshqaruv moduli")]
    sw = 2.92
    for i, (n, d) in enumerate(stats):
        x = 0.7 + i * (sw + 0.13)
        rrect(s, x, 2.75, sw, 1.35, fill=DEEP, radius=0.12)
        txt(s, n, x, 2.82, sw, 0.7, size=40, color=GOLD, bold=True, font=HEAD,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        txt(s, d, x + 0.15, 3.55, sw - 0.3, 0.5, size=11.5, color="CFE6E2", bold=True,
            font=BODY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP, margin=0.0, spacing=1.0)
    # natija bulletlari
    points = [
        "Qo'lda yuritish o'rniga — yagona, avtomatlashtirilgan tizim.",
        "Sun'iy intellekt: churn bashorati, kurs tavsiyasi, o'qituvchi hisoboti.",
        "Veb + mobil + Telegram botlar — barcha ishtirokchilar uchun qulay.",
        "Coin va reyting tizimi orqali talabalar motivatsiyasi oshiriladi.",
    ]
    yy = 4.4
    for p in points:
        oval(s, 0.78, yy + 0.06, 0.18, 0.18, fill=GOLD)
        txt(s, p, 1.15, yy - 0.06, 11.0, 0.45, size=14, color="E7F2F0", font=BODY,
            valign=MSO_ANCHOR.MIDDLE, margin=0.0)
        yy += 0.5
    # rahmat
    rrect(s, 0.7, 6.45, 11.95, 0.6, fill=BRAND, radius=0.3)
    txt(s, "E'tiboringiz uchun rahmat!  ·  Savollarga tayyorman", 0.7, 6.45, 11.95, 0.6,
        size=15, color=WHITE, bold=True, font=HEAD, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE, margin=0.0)


# ── Qurish ────────────────────────────────────────────────
slide1()
slide2()
slide3()
slide4()
slide5()
slide6()
# 7 — Dashboard
feature_slide(7, "05", "Boshqaruv paneli (Dashboard)", "REAL VAQT KO'RSATKICHLARI",
              "03_dashboard.png",
              [("Jonli statistika", "Talaba, kurs, daromad va davomat bir ekranda."),
               ("Moliyaviy ko'rsatkich", "Oylik daromad va kutilayotgan to'lovlar."),
               ("Tezkor navigatsiya", "Barcha modullarga bitta paneldan kirish.")],
              "Boshqaruv paneli — markazning umumiy holati bir qarashda", color=BRAND,
              advantage="Yagona panelda biznesning barcha asosiy raqamlari real vaqtda.")
# 8 — Talabalar
feature_slide(8, "06", "Talabalar boshqaruvi", "STUDENTS",
              "04_students.png",
              [("To'liq profil", "Telefon, ota-ona, jinsi, holati va rasm."),
               ("Qidiruv va filtr", "Faol/nofaol bo'yicha tez topish."),
               ("Ro'yxatga ulash", "Talabani bir qadamda kursga yozish.")],
              "Talabalar ro'yxati — qidiruv, filtr va holat boshqaruvi", color=BRAND,
              advantage="Har bir talaba uchun yagona, to'liq elektron kartoteka.")
# 9 — AI churn
ai_slide(9, "07", "AI — talaba churn (ketib qolish) bashorati",
         "18_ai_churn_students.png",
         [("Ma'lumot yig'ish", "Davomat, to'lov va kurs yozilishlari tahlil qilinadi."),
          ("AI baholash", "Claude AI 0–100 ball va o'zbekcha sabab beradi."),
          ("Avtomatik ogohlantirish", "Ball 70 dan oshsa — adminlarga signal yuboriladi.")],
         "Talabaning ketib qolish ehtimolini oldindan bashorat qiladi va xavf yuqori "
         "bo'lganda avtomatik ogohlantiradi — talabani yo'qotmasdan oldin harakat qilish imkoni.",
         "Talabalar ro'yxatidagi «AI Risk» ustuni — rangli xavf darajalari")
# 10 — Kurslar + transfer
feature_slide(10, "08", "Kurslar va talaba transferi", "COURSES",
              "06_courses.png",
              [("Kurslar katalogi", "Daraja, narx, davomiylik va o'qituvchi."),
               ("Talaba transferi", "Bir kursdan boshqasiga bir bosishda o'tkazish."),
               ("Yangi kurs", "Modal oynada tez kurs yaratish.")],
              "Kurslar bo'limi — katalog va boshqaruv", color=TEAL,
              second=("06b_transfer_modal.png", "«Talabani boshqa kursga o'tkazish» oynasi"),
              advantage="Talabani kursdan-kursga o'tkazish — qo'lda emas, bir necha soniyada.")
# 11 — Ro'yxatga olish
feature_slide(11, "09", "Ro'yxatga olish (Enrollments)", "ENROLLMENTS",
              "07_enrollments.png",
              [("Yozilish boshqaruvi", "Talaba–kurs bog'lanishi va statusi."),
               ("To'lov holati", "To'langan/qarz foizi har bir yozilishda."),
               ("Tez amallar", "Tahrirlash, status va o'chirish.")],
              "Ro'yxatga olish — talaba va kurs bog'lanishi", color=BRAND,
              advantage="Ro'yxatga olish va to'lov holati bitta joyda kuzatiladi.")
# 12 — AI kurs tavsiyasi
ai_slide(12, "10", "AI — kurs tavsiyasi",
         "20_ai_course_recommend.png",
         [("Tarixni tahlil", "Talabaning tugatgan va faol kurslari o'rganiladi."),
          ("Mos kurs tanlash", "Claude AI mavjud kurslardan eng mos 3 tasini tanlaydi."),
          ("Asoslangan tavsiya", "Har bir tavsiyaga o'zbekcha sabab beriladi.")],
         "Talabaga keyingi bosqich uchun eng mos kurslarni AI tavsiya qiladi — "
         "sotuvni oshiradi va talabaning o'qish yo'lini izchil davom ettiradi.",
         "Ro'yxatga olishda «AI kurs tavsiyasi» — 3 ta asoslangan tavsiya")
# 13 — O'qituvchilar
feature_slide(13, "11", "O'qituvchilar boshqaruvi", "INSTRUCTORS",
              "16_instructors.png",
              [("O'qituvchi va admin", "Rollar va ruxsatlar boshqaruvi."),
               ("Tajriba va indeks", "Mutaxassislik, tajriba va ko'rsatkichlar."),
               ("Holat nazorati", "Faol/nofaol o'qituvchilarni kuzatish.")],
              "O'qituvchilar bo'limi — kadrlar boshqaruvi", color=BRAND,
              advantage="Barcha kadrlar va ularning ko'rsatkichlari yagona ro'yxatda.")
# 14 — AI o'qituvchi hisoboti
ai_slide(14, "12", "AI — o'qituvchi samaradorlik hisoboti",
         "19_ai_instructor_report.png",
         [("Statistika yig'ish", "Retention, davomat va yozilishlar to'planadi."),
          ("AI hisobot", "Claude AI o'zbek tilida 3–5 jumlali xulosa yozadi."),
          ("Kuchli/zaif tomon", "Yaxshilanishi kerak bo'lgan jihatlar ko'rsatiladi.")],
         "O'qituvchi samaradorligini raqamlar asosida obyektiv baholaydi va tushunarli "
         "tilda hisobot yozadi — kadrlar bo'yicha qaror qabul qilishni osonlashtiradi.",
         "O'qituvchilarda «AI samaradorlik hisoboti» — obyektiv tahlil")
slide15()
slide16()
slide17()
# 18 — O'qituvchi mobil ilova
mobile_slide(18, "16", "O'qituvchi mobil ilovasi", "Faqat o'qituvchilar uchun — JWT autentifikatsiya bilan xavfsiz kirish.",
             BRAND,
             [("To'lovlar", "Guruh to'lovlarini ko'rish va boshqarish."),
              ("Guruhlar", "O'z guruhlari va talabalar ro'yxati."),
              ("Davomat olish", "Darsda davomatni telefondan belgilash.")],
             ["To'lovlar", "Guruhlar", "Davomat"])
# 19 — Talaba mobil ilova
mobile_slide(19, "17", "Talaba mobil ilovasi", "Talabalar uchun — to'lov, motivatsiya va dars jadvali bir ilovada.",
             TEAL,
             [("To'lov qilish", "Kurs to'lovini ilova orqali amalga oshirish."),
              ("Coin & reyting", "To'plangan coinlar va reyting o'rni."),
              ("Dars kunlari", "Dars jadvali va kelgusi mashg'ulotlar.")],
             ["To'lov qilish", "Coinlar", "Reyting", "Dars kunlari"])
slide20()

prs.save(OUT)
print("SAQLANDI:", OUT)
print("Slaydlar:", len(prs.slides._sldIdLst))

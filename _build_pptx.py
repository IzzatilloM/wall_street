# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image

SHOTS = os.path.join(os.path.dirname(__file__), "_pptx_shots")
OUT = os.path.join(os.path.dirname(__file__), "Wall_Street_Tizim_Prezentatsiya.pptx")

# ---------- palette ----------
BG    = RGBColor(0x0B, 0x2A, 0x2C)
BG2   = RGBColor(0x0E, 0x35, 0x38)
CARD  = RGBColor(0x10, 0x3E, 0x42)
CARD2 = RGBColor(0x14, 0x4A, 0x4E)
TEAL  = RGBColor(0x0E, 0x5C, 0x60)
MINT  = RGBColor(0x2E, 0xC4, 0xB6)
GOLD  = RGBColor(0xE0, 0xB8, 0x4C)
TXT   = RGBColor(0xEC, 0xF4, 0xF3)
MUT   = RGBColor(0x9C, 0xB7, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HFONT = "Trebuchet MS"
BFONT = "Calibri"

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=BG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def _set_font(run, size, color, bold=False, italic=False, font=BFONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line=None, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic,font)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line:
            p.line_spacing = line
        for seg in para:
            t, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            ital = seg[4] if len(seg) > 4 else False
            fnt  = seg[5] if len(seg) > 5 else BFONT
            r = p.add_run(); r.text = t
            _set_font(r, size, color, bold, ital, fnt)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'),
              {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val':'000000'})
        alpha = clr.makeelement(qn('a:alpha'), {'val':'45000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def img_fit(path, box_w, box_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    bar = box_w / box_h
    if ar > bar:
        w = box_w; h = box_w / ar
    else:
        h = box_h; w = box_h * ar
    return w, h

def framed_image(s, path, cx, cy, box_w, box_h, border=GOLD, bw=1.5):
    """Place image fit within box, centered at (cx,cy) area top-left = cx,cy box."""
    w, h = img_fit(path, box_w, box_h)
    x = cx + (box_w - w) / 2
    y = cy + (box_h - h) / 2
    # shadow backing
    rect(s, x, y, w, h, fill=RGBColor(0,0,0), shadow=True)
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    # border
    rect(s, x, y, w, h, fill=None, line=border, line_w=bw)
    return x, y, w, h

def num_badge(s, x, y, n, d=0.62):
    rect(s, x, y, d, d, fill=GOLD, shape=MSO_SHAPE.OVAL)
    text(s, x, y-0.02, d, d, [[(str(n), 24, BG, True, False, HFONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def bullets(s, x, y, w, items, gap=0.62, dot=GOLD, head_size=13.5, body_size=11):
    cy = y
    for head, body in items:
        rect(s, x, cy+0.07, 0.14, 0.14, fill=dot, shape=MSO_SHAPE.OVAL)
        text(s, x+0.30, cy, w-0.30,
             0.9, [[(head+"  ", head_size, TXT, True, False, HFONT),
                    (body, body_size, MUT, False, False, BFONT)]],
             line=1.04, space_after=0)
        cy += gap
    return cy

def content_slide(n, kicker, title, intro, items, img, img_side="right",
                  intro_color=MUT, accent=GOLD):
    s = slide(); bg(s, BG if n % 2 else BG2)
    # decorative corner dots motif (subtle)
    rect(s, SW-0.55, 0.32, 0.16, 0.16, fill=TEAL, shape=MSO_SHAPE.OVAL)
    rect(s, SW-0.32, 0.32, 0.16, 0.16, fill=GOLD, shape=MSO_SHAPE.OVAL)

    text_w = 4.7
    if img_side == "right":
        tx = 0.7
        img_x = 5.75
    else:
        tx = SW - text_w - 0.7
        img_x = 0.7
    img_w = SW - 5.75 - 0.7 if img_side == "right" else 5.75
    img_w = 6.85
    img_box_x = img_x

    num_badge(s, tx, 0.5, n)
    text(s, tx+0.85, 0.52, text_w-0.5, 0.4,
         [[(kicker.upper(), 12, MINT, True, False, HFONT)]], space_after=0)
    text(s, tx, 1.32, text_w, 1.2,
         [[(title, 26, TXT, True, False, HFONT)]], line=0.98, space_after=0)
    text(s, tx, 2.78, text_w, 1.2,
         [[(intro, 12.5, intro_color, False, False, BFONT)]], line=1.12, space_after=0)
    bullets(s, tx, 4.05, text_w, items)

    framed_image(s, img, img_box_x, 1.15, img_w, 5.2, border=accent, bw=1.5)
    return s

# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = slide(); bg(s, BG)
# big soft circles motif
rect(s, 9.4, -1.8, 6.0, 6.0, fill=BG2, shape=MSO_SHAPE.OVAL)
rect(s, 10.8, 3.6, 4.4, 4.4, fill=CARD, shape=MSO_SHAPE.OVAL)
# logo badge
rect(s, 0.9, 0.85, 1.25, 1.25, fill=TEAL, shape=MSO_SHAPE.OVAL, line=GOLD, line_w=2.0)
text(s, 0.9, 0.9, 1.25, 1.15, [[("WS", 30, GOLD, True, False, HFONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 2.35, 1.02, 7, 0.5, [[("WALL STREET", 22, TXT, True, False, HFONT)]], space_after=0)
text(s, 2.36, 1.55, 7, 0.4, [[("EDUCATION CENTRE", 13, MINT, True, False, HFONT)]], space_after=0)

text(s, 0.9, 2.85, 9.6, 2.0,
     [[("O'quv markazini boshqarish", 46, TXT, True, False, HFONT)],
      [("to'liq raqamli tizimi", 46, GOLD, True, False, HFONT)]],
     line=1.0, space_after=2)
text(s, 0.95, 4.75, 8.6, 0.8,
     [[("Web CRM · Mobil ilova · Telegram bot · REST API — talaba, kurs, "
        "to'lov, davomat va coin tizimi yagona platformada.", 15, MUT, False, False, BFONT)]],
     line=1.2)
# bottom chips
chips = ["51 talaba", "21 kurs", "22 o'qituvchi", "Coin tizimi", "PDF cheklar"]
cx = 0.95
for c in chips:
    w = 0.32 + len(c)*0.105
    rect(s, cx, 5.85, w, 0.46, fill=CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5, line=TEAL, line_w=1.0)
    text(s, cx, 5.85, w, 0.46, [[(c, 11.5, TXT, True, False, BFONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + 0.18
text(s, 0.95, 6.75, 9, 0.4, [[("Tizim to'liq tavsifi · 2026", 11, MUT, False, True, BFONT)]])

# ============================================================
# SLIDE 2 — ARCHITECTURE / COMPONENTS
# ============================================================
s = slide(); bg(s, BG2)
text(s, 0.7, 0.55, 11, 0.5, [[("TIZIM ARXITEKTURASI", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.0, 12, 0.8, [[("To'rtta bog'langan komponent, bitta baza", 30, TXT, True, False, HFONT)]], space_after=0)
comps = [
    ("Web CRM (Django)", "Administrator va o'qituvchilar uchun to'liq boshqaruv paneli: "
     "talabalar, kurslar, davomat, to'lovlar, hisobotlar.", GOLD),
    ("Mobil ilova (Flutter)", "Talaba o'z profili, kurslari, to'lovlari va coin balansini "
     "telefonida ko'radi. JWT token orqali xavfsiz kirish.", MINT),
    ("Telegram bot (aiogram)", "Yangi talabalarni ro'yxatdan o'tkazadi, lead yig'adi va "
     "kurslar haqida ma'lumot beradi — webhook orqali sayt bilan bitta jarayonda.", GOLD),
    ("REST API (DRF + JWT)", "Mobil ilova va saytni bog'lovchi xavfsiz interfeys: "
     "auth, profil, kurslar, to'lovlar va yangiliklar endpointlari.", MINT),
]
gx, gy, gw, gh = 0.7, 2.1, 5.85, 1.95
positions = [(0.7,2.1),(6.78,2.1),(0.7,4.25),(6.78,4.25)]
for (px,py),(h,b,acc) in zip(positions, comps):
    rect(s, px, py, gw, gh, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, line=TEAL, line_w=1.0)
    rect(s, px, py, 0.12, gh, fill=acc, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, px+0.4, py+0.3, gw-0.7, 0.5, [[(h, 17, TXT, True, False, HFONT)]], space_after=0)
    text(s, px+0.4, py+0.85, gw-0.7, 1.0, [[(b, 12, MUT, False, False, BFONT)]], line=1.12)

# ============================================================
# CONTENT SLIDES 3..18
# ============================================================
content_slide(3, "Xavfsiz kirish",
    "Kirish va ro'yxatdan o'tish",
    "Tizimga username yoki Gmail orqali kiriladi. Har bir foydalanuvchi roli (admin / "
    "o'qituvchi / talaba) bo'yicha o'z paneliga yo'naltiriladi.",
    [("Rolga asoslangan kirish", "— admin to'liq, o'qituvchi cheklangan huquq."),
     ("Parolni ko'rsatish", "va \"meni eslab qol\" imkoniyati."),
     ("Parolni tiklash", "— SMS kod orqali tasdiqlash."),
     ("Telegram bot", "orqali avtomatik login/parol berish.")],
    os.path.join(SHOTS,"01_login.png"), "right")

content_slide(4, "Boshqaruv paneli",
    "Dashboard — bir qarashda butun markaz",
    "Bosh sahifada markazning asosiy ko'rsatkichlari jonli ravishda chiqadi: oylik "
    "daromad, kutilayotgan to'lovlar, faol talabalar va kurslar.",
    [("Moliyaviy kartalar", "— oylik daromad va qarzlar."),
     ("Statistika", "— talabalar, kurslar, o'qituvchilar soni."),
     ("Davomat xulosasi", "— bugungi darslar holati."),
     ("Tez navigatsiya", "— chap menyudan barcha bo'limlarga.")],
    os.path.join(SHOTS,"03_dashboard.png"), "left")

content_slide(5, "Talabalar",
    "Talabalar ro'yxati va boshqaruvi",
    "Barcha talabalar yagona ro'yxatda. Qidiruv, filtrlash, yangi talaba qo'shish va "
    "har bir talabaning holatini (faol/nofaol) boshqarish.",
    [("Qidiruv va filtr", "— ism, ID yoki holat bo'yicha."),
     ("Yangi talaba", "qo'shish va tahrirlash."),
     ("Holat almashtirish", "— faol / nofaol qilish."),
     ("Profilga o'tish", "— bitta bosishda to'liq ma'lumot.")],
    os.path.join(SHOTS,"04_students.png"), "right")

content_slide(6, "Talaba profili",
    "Har bir talaba — to'liq kartochka",
    "Talaba profilida shaxsiy ma'lumotlar, yozilgan kurslari, to'lov tarixi, davomat "
    "foizi va coin balansi birgalikda ko'rinadi.",
    [("Shaxsiy ma'lumot", "— telefon, ota-ona, manzil, rasm."),
     ("Kurslari va to'lovlari", "bitta joyda."),
     ("Davomat foizi", "har bir kurs bo'yicha."),
     ("Coin balansi", "— rag'batlantirish ko'rsatkichi.")],
    os.path.join(SHOTS,"05_student_detail.png"), "left")

content_slide(7, "Kurslar",
    "Kurslarni boshqarish paneli",
    "Markazdagi barcha kurslar daraja va narx bilan. Kursni yaratish, tahrirlash, "
    "o'chirish va ichidagi talabalarni ko'rish.",
    [("Kurslar katalogi", "— daraja, narx, davomiylik."),
     ("Qidiruv va filtr", "daraja bo'yicha."),
     ("Kurs ichidagilar", "— talabalar ro'yxati."),
     ("Statistika", "— jami kurs, talaba, o'qituvchi.")],
    os.path.join(SHOTS,"06_courses.png"), "right")

content_slide(8, "Yangi kurs",
    "Kurs qo'shish va jadval",
    "Yangi kursni yaratishda nomi, darajasi, narxi, o'qituvchisi va dars jadvali "
    "(hafta kunlari, vaqt) belgilanadi — barchasi bitta oynada.",
    [("Modal oyna", "— sahifani tark etmasdan yaratish."),
     ("Daraja va narx", "belgilash."),
     ("O'qituvchi biriktirish", "kursga."),
     ("Dars jadvali", "— kun va vaqt sozlamalari.")],
    os.path.join(SHOTS,"06c_create_course_modal.png"), "left")

# SLIDE 9 — TRANSFER (highlighted)
content_slide(9, "Maxsus imkoniyat",
    "Talabani boshqa kursga ko'chirish",
    "Bitta bosish bilan talabani bir kursdan ikkinchisiga o'tkazish mumkin. Modal oynada "
    "talaba ismi va hozirgi kursi ko'rsatiladi, yangi kurs ro'yxatdan tanlanadi.",
    [("Transfer tugmasi", "— har bir talaba qatorida."),
     ("Joriy kurs", "avtomatik ko'rsatiladi."),
     ("Yangi kurs tanlash", "— ochiluvchi ro'yxat."),
     ("Tezkor o'tkazish", "— ma'lumot yo'qolmaydi.")],
    os.path.join(SHOTS,"06b_transfer_modal.png"), "right")

content_slide(10, "Ro'yxatga olish",
    "Enrollment — talabani kursga yozish",
    "Talaba qaysi kursga, qaysi o'qituvchiga yozilgani, to'lov holati va foizi shu "
    "bo'limda boshqariladi. Telegram leadlari ham shu yerga tushadi.",
    [("Kursga yozish", "— talaba + kurs + o'qituvchi."),
     ("To'lov holati", "— to'langan / qisman / qarz."),
     ("Holatni almashtirish", "va izoh qo'shish."),
     ("Telegram leadlar", "— botdan kelganlar.")],
    os.path.join(SHOTS,"07_enrollments.png"), "left")

content_slide(11, "Davomat",
    "Attendance — darslar davomati",
    "O'qituvchi har bir dars uchun davomat belgilaydi: kelgan, kelmagan, kechikkan. "
    "Tizim foizlarni avtomatik hisoblaydi va CSV ga eksport qiladi.",
    [("Davomat belgilash", "— present / absent / late."),
     ("Avtomatik foiz", "har bir talaba bo'yicha."),
     ("Filtrlar", "— kurs, sana, holat."),
     ("CSV eksport", "— hisobot uchun.")],
    os.path.join(SHOTS,"08_attendance.png"), "right")

# SLIDE 12 — COIN (highlighted)
content_slide(12, "Rag'batlantirish",
    "Coin tizimi — talabani motivatsiya qilish",
    "O'qituvchi davomat oynasidan talabaga coin qo'shadi yoki ayiradi (faollik, "
    "uy vazifasi uchun). Har bir o'zgarish tarix sifatida saqlanadi va talaba "
    "balansiga ta'sir qiladi.",
    [("\"Coin berish\"", "— bir tugma orqali."),
     ("Qo'shish / ayirish", "— sababi bilan."),
     ("Tarix saqlanadi", "— kim, qachon, nega."),
     ("Mobil ilovada", "balans talabaga ko'rinadi.")],
    os.path.join(SHOTS,"09_coin_award.png"), "left")

content_slide(13, "To'lovlar",
    "To'lovlar ro'yxati va nazorati",
    "Barcha to'lovlar yagona ro'yxatda: summa, holat (to'langan / kutilmoqda), oy va "
    "talaba. Har bir to'lovni ko'rish, tahrirlash va chek chiqarish mumkin.",
    [("To'lovlar ro'yxati", "— summa va holat."),
     ("Filtr va qidiruv", "— oy, talaba bo'yicha."),
     ("Holatlar", "— to'langan / pending."),
     ("Tezkor amallar", "— ko'rish, tahrir, o'chirish.")],
    os.path.join(SHOTS,"10_payments.png"), "right")

content_slide(14, "To'lov qabul qilish",
    "Yangi to'lov qo'shish",
    "To'lovni qabul qilishda talaba va kurs tanlanadi, narx avtomatik chiqadi, to'lov "
    "usuli (naqd / karta / o'tkazma) va oy belgilanadi.",
    [("Talaba + kurs", "tanlash."),
     ("Avtomatik narx", "— kurs bo'yicha."),
     ("To'lov usuli", "— naqd / karta / o'tkazma."),
     ("Oy belgilash", "— qaysi oy uchun.")],
    os.path.join(SHOTS,"11_payment_create.png"), "left")

content_slide(15, "Chek va kvitansiya",
    "To'lov cheki — PDF, QR va pechat",
    "Har bir to'lov uchun raqamli kvitansiya: talaba, kurs, o'qituvchi va summa. "
    "Bir tugma bilan QR-kod va markaz pechati bosilgan PDF chek yuklab olinadi.",
    [("To'liq tafsilot", "— talaba, kurs, oy, summa."),
     ("PDF chek", "— yuklab olish / chop etish."),
     ("QR-kod", "— tekshirish uchun."),
     ("Pechat va raqam", "— rasmiy kvitansiya.")],
    os.path.join(SHOTS,"12_payment_detail.png"), "right")

content_slide(16, "Hisobotlar",
    "Reports — moliyaviy tahlil",
    "Markaz faoliyatining umumiy tahlili: daromad dinamikasi, to'lovlar, talabalar "
    "o'sishi va kurslar bo'yicha taqsimot grafiklarda.",
    [("Daromad tahlili", "— davr bo'yicha."),
     ("Grafik va diagrammalar", "— vizual ko'rinish."),
     ("To'lov statistikasi", "— to'langan / qarz."),
     ("Qaror qabul qilish", "uchun jonli ma'lumot.")],
    os.path.join(SHOTS,"13_reports.png"), "left")

content_slide(17, "Xodimlar maoshi",
    "Salary — maoshlarni boshqarish",
    "O'qituvchi va xodimlar maoshlari hisoblanadi va saqlanadi. Maosh ro'yxatini "
    "Excel va PDF formatida eksport qilish mumkin.",
    [("Maosh yozuvlari", "— xodim bo'yicha."),
     ("Qo'shish va tahrir", "qilish."),
     ("Excel eksport", "— hisob-kitob uchun."),
     ("PDF eksport", "— rasmiy hujjat.")],
    os.path.join(SHOTS,"14_salary.png"), "right")

# ============================================================
# SLIDE 18 — Calendar + Instructors + Settings (triple)
# ============================================================
s = slide(); bg(s, BG2)
num_badge(s, 0.7, 0.55, 18)
text(s, 1.55, 0.56, 9, 0.4, [[("QO'SHIMCHA BO'LIMLAR", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 12, 0.7, [[("Kalendar · O'qituvchilar · Sozlamalar", 28, TXT, True, False, HFONT)]], space_after=0)
trip = [
    (os.path.join(SHOTS,"15_calendar.png"), "Kalendar", "Darslar, imtihonlar va tadbirlar oylik kalendarda; kunlar bo'yicha tafsilot."),
    (os.path.join(SHOTS,"16_instructors.png"), "O'qituvchilar", "O'qituvchilar ro'yxati, ularning kurslari va yuklamasi bir joyda."),
    (os.path.join(SHOTS,"17_settings.png"), "Sozlamalar", "Markaz nomi, logo, narxlar va tizim parametrlarini moslash."),
]
cw = 3.95; gap = 0.28; x0 = 0.7; iy = 2.2
for i,(p,h,b) in enumerate(trip):
    px = x0 + i*(cw+gap)
    rect(s, px, iy, cw, 3.9, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, line=TEAL, line_w=1.0)
    framed_image(s, p, px+0.18, iy+0.18, cw-0.36, 2.35, border=GOLD, bw=1.0)
    text(s, px+0.28, iy+2.75, cw-0.56, 0.4, [[(h, 16, GOLD, True, False, HFONT)]], space_after=0)
    text(s, px+0.28, iy+3.18, cw-0.56, 0.7, [[(b, 11, MUT, False, False, BFONT)]], line=1.1)

# ============================================================
# SLIDE 19 — TELEGRAM BOT
# ============================================================
s = slide(); bg(s, BG)
num_badge(s, 0.7, 0.55, 19)
text(s, 1.55, 0.56, 9, 0.4, [[("TELEGRAM BOT", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 7, 1.0, [[("Bot orqali ro'yxatdan o'tish", 30, TXT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 2.15, 6.0, 1.2,
     [[("Telegram bot yangi talabalarni avtomatik ro'yxatga oladi va lead yig'adi. "
        "Webhook orqali ishlaydi — sayt bilan bitta jarayon, bitta baza.", 13.5, MUT, False, False, BFONT)]],
     line=1.18)
bot_items = [
    ("/start menyusi", "— kurslar, narxlar, manzil, ro'yxatdan o'tish."),
    ("Telefon va ism", "so'raydi, FSM holatni bazada saqlaydi."),
    ("Avtomatik login", "— tizimga kirish ma'lumotlarini beradi."),
    ("Lead yig'ish", "— qiziqqanlar CRM enrollmentlariga tushadi."),
    ("Kurslar ro'yxati", "— inline tugmalar bilan tanlash."),
]
bullets(s, 0.7, 3.5, 6.0, [(h,b) for h,b in bot_items], gap=0.6, head_size=13.5, body_size=11.5)

# phone mock with chat on the right
pmx, pmy, pmw, pmh = 8.0, 1.15, 4.4, 5.9
rect(s, pmx, pmy, pmw, pmh, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, line=GOLD, line_w=1.5, shadow=True)
# header
rect(s, pmx, pmy, pmw, 0.95, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
rect(s, pmx, pmy+0.45, pmw, 0.5, fill=TEAL)
rect(s, pmx+0.28, pmy+0.22, 0.5, 0.5, fill=GOLD, shape=MSO_SHAPE.OVAL)
text(s, pmx+0.30, pmy+0.27, 0.5, 0.4, [[("WS", 12, BG, True, False, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, pmx+0.95, pmy+0.20, 3, 0.3, [[("Wall Street bot", 13, WHITE, True, False, HFONT)]], space_after=0)
text(s, pmx+0.95, pmy+0.52, 3, 0.3, [[("online", 10, RGBColor(0xCD,0xEA,0xE8), False, False, BFONT)]], space_after=0)
# chat bubbles
def bubble(y, who, msg, h):
    if who == "bot":
        rect(s, pmx+0.25, y, 3.0, h, fill=CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        text(s, pmx+0.42, y+0.12, 2.7, h-0.2, [[(msg, 10.5, TXT, False, False, BFONT)]], line=1.08, space_after=0)
    else:
        rect(s, pmx+1.15, y, 3.0, h, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        text(s, pmx+1.32, y+0.12, 2.7, h-0.2, [[(msg, 10.5, WHITE, False, False, BFONT)]], line=1.08, space_after=0)
bubble(pmy+1.15, "bot", "Assalomu alaykum! Wall Street'ga xush kelibsiz 👋\nRo'yxatdan o'tishni boshlaymizmi?", 0.95)
bubble(pmy+2.25, "user", "Ha, ro'yxatdan o'taman ✅", 0.5)
bubble(pmy+2.90, "bot", "📱 Telefon raqamingizni yuboring:", 0.5)
bubble(pmy+3.55, "user", "+998 90 123 45 67", 0.5)
bubble(pmy+4.20, "bot", "Tayyor! Login: ali.valiyev\nParol: ••••••••\nIlovaga shu bilan kiring 🎉", 0.95)

# ============================================================
# SLIDE 20 — MOBILE APP + CLOSING
# ============================================================
s = slide(); bg(s, BG2)
num_badge(s, 0.7, 0.55, 20)
text(s, 1.55, 0.56, 9, 0.4, [[("MOBIL ILOVA (FLUTTER)", 12, MINT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 1.2, 7, 1.0, [[("Talaba uchun mobil ilova", 30, TXT, True, False, HFONT)]], space_after=0)
text(s, 0.7, 2.15, 6.0, 1.0,
     [[("Talaba o'qish jarayonini telefonida kuzatadi. REST API (JWT) orqali "
        "tizim bilan xavfsiz bog'lanadi.", 13.5, MUT, False, False, BFONT)]], line=1.18)
mob_items = [
    ("Profil", "— shaxsiy ma'lumotlarni ko'rish va yangilash."),
    ("Mening kurslarim", "— daraja, o'qituvchi, to'lov foizi."),
    ("To'lovlar", "— to'langan, qarz va umumiy xulosa."),
    ("Coin balansi", "— rag'batlantirish ballari."),
    ("Yangiliklar", "— markaz e'lonlari."),
]
bullets(s, 0.7, 3.35, 6.0, [(h,b) for h,b in mob_items], gap=0.55, head_size=13.5, body_size=11.5)

# phones
framed_image(s, os.path.join(SHOTS,"M1_login.png"), 7.7, 1.15, 2.55, 5.95, border=GOLD, bw=1.5)
framed_image(s, os.path.join(SHOTS,"M3_after.png"), 10.35, 1.15, 2.55, 5.95, border=TEAL, bw=1.5)
text(s, 7.7, 7.05, 5.2, 0.3,
     [[("REST API: /api/auth · /api/me · /api/me/courses · /api/me/payments · /api/news",
        9.5, MUT, False, True, BFONT)]], align=PP_ALIGN.CENTER, space_after=0)

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))
